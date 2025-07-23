import os
import logging
from typing import Optional, Dict, Any
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import magic

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ContentExtractor:
    """Extracts text content from various file types."""
    
    def __init__(self):
        self.mime = magic.Magic(mime=True)
    
    def extract_content(self, file_path: str, file_meta: Optional[Dict] = None) -> Dict[str, Any]:
        """Extract content from a file based on its type.
        
        Args:
            file_path: Path to the file
            file_meta: Optional file metadata (if already available)
            
        Returns:
            Dictionary with extracted content and metadata
        """
        if file_meta is None:
            from .file_processor import FileProcessor
            file_meta = FileProcessor().get_file_metadata(file_path)
            
        if not file_meta:
            return {}
        
        content = {
            'path': file_meta['path'],
            'file_type': file_meta.get('file_type', 'unknown'),
            'content': '',
            'content_type': file_meta.get('content_type', 'unknown'),
            'error': None
        }
        
        try:
            if file_meta['extension'] in ['.pdf']:
                content['content'] = self._extract_pdf_text(file_path)
            elif file_meta['extension'] in ['.docx', '.doc']:
                content['content'] = self._extract_docx_text(file_path)
            elif file_meta['extension'] in ['.xlsx', '.xls']:
                content['content'] = self._extract_excel_text(file_path)
            elif file_meta['extension'] in ['.pptx', '.ppt']:
                content['content'] = self._extract_pptx_text(file_path)
            elif file_meta['extension'] in ['.html', '.htm']:
                content['content'] = self._extract_html_text(file_path)
            elif file_meta['extension'] in ['.txt', '.md', '.csv', '.json']:
                content['content'] = self._extract_plain_text(file_path)
            elif file_meta['content_type'] == 'image':
                # For images, we'll handle features separately
                content['content'] = ''
            else:
                logger.warning(f"Unsupported file type for content extraction: {file_path}")
                content['error'] = 'Unsupported file type'
                
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {str(e)}")
            content['error'] = str(e)
        
        return content
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files."""
        text = []
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text.append(page.extract_text() or '')
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            return ''
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from Word documents."""
        try:
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
            return ''
    
    def _extract_excel_text(self, file_path: str) -> str:
        """Extract text from Excel files."""
        try:
            text_parts = []
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                text_parts.append(df.to_string())
            return '\n'.join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {str(e)}")
            return ''
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations."""
        try:
            prs = Presentation(file_path)
            text = []
            for i, slide in enumerate(prs.slides):
                text.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
            return ''
    
    def _extract_html_text(self, file_path: str) -> str:
        """Extract text from HTML files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return ' '.join(soup.stripped_strings)
        except Exception as e:
            logger.error(f"Error extracting text from HTML {file_path}: {str(e)}")
            return ''
    
    def _extract_plain_text(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {str(e)}")
            return ''

# Example usage
if __name__ == "__main__":
    extractor = ContentExtractor()
    test_file = "path/to/test/file.docx"
    if os.path.exists(test_file):
        content = extractor.extract_content(test_file)
        print(f"Extracted content from {test_file}:")
        print(content['content'][:500] + "..." if len(content['content']) > 500 else content['content'])
    else:
        print(f"Test file not found: {test_file}")
